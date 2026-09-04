"""
AURA :: Document Builder
========================
Turns a structured outline into a REAL Office file: .pptx, .xlsx or .docx.

WHAT THIS MODULE IS, AND IS NOT
-------------------------------
It is the *rendering* half only. It takes an already-structured outline —
a dict of slides / rows / sections — and writes a genuine OOXML file to disk.

It deliberately contains **no AI**. The prompt-to-outline step happens in the
browser (`js/ai/doc-agent.js`) using whichever provider the user already
configured, and its JSON is validated here before a single byte is written.
That split matters: the dangerous part (writing files) stays small, auditable
and fully testable without a model running.

HONESTY
-------
python-pptx / openpyxl / python-docx are OPTIONAL. If one is missing, the
matching builder reports exactly which `pip install` fixes it and refuses.
It never writes a fake file, never writes a .txt and calls it a .pptx.

SAFETY
------
Every output path goes through `bridge._resolve_path`, so the same jail that
protects the file plugin protects this one: no writing outside the user's
folders, no symlink escape, no credential directories. The extension is forced
to match the builder, so a "presentation" can never be written as `.exe`.
"""

import os
import re
import datetime

# The builders are optional; report per-format rather than failing at import.
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except Exception:                                  # pragma: no cover
    HAS_PPTX = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter
    HAS_XLSX = True
except Exception:                                  # pragma: no cover
    HAS_XLSX = False

try:
    import docx
    from docx.shared import Pt as DocxPt, RGBColor as DocxRGB
    HAS_DOCX = True
except Exception:                                  # pragma: no cover
    HAS_DOCX = False


MAX_SLIDES = 40
MAX_ROWS = 5000
MAX_COLS = 60
MAX_SECTIONS = 80
MAX_TEXT = 4000

# AURA's accent, so generated decks look like they came from AURA.
ACCENT = (0x38, 0xBD, 0xF8)
INK = (0x10, 0x18, 0x24)
DIM = (0x5A, 0x6B, 0x80)


def capabilities():
    """What can actually be built on this machine, right now."""
    return {
        "ok": True,
        "pptx": HAS_PPTX,
        "xlsx": HAS_XLSX,
        "docx": HAS_DOCX,
        "any": HAS_PPTX or HAS_XLSX or HAS_DOCX,
        "install": {
            "pptx": "pip install python-pptx",
            "xlsx": "pip install openpyxl",
            "docx": "pip install python-docx",
        },
        "note": ("Each format needs its own library. AURA reports which one is "
                 "missing instead of writing a placeholder file."),
    }


def _slug(name, fallback="aura-document"):
    """A safe file stem. Never a path — separators are stripped, not escaped."""
    s = re.sub(r"[^\w\s.-]", "", str(name or "")).strip()
    s = re.sub(r"[\s_]+", "-", s).strip("-.")
    s = s[:60] or fallback
    # Windows reserved device names would make an unopenable file.
    if s.upper().split(".")[0] in {
            "CON", "PRN", "AUX", "NUL", "COM1", "COM2", "COM3", "COM4",
            "LPT1", "LPT2", "LPT3"}:
        s = f"aura-{s}"
    return s


def _target_path(folder, name, ext, resolver):
    """
    Build and jail-check the output path.

    `resolver` is injected (bridge._resolve_path) so this module stays testable
    without importing bridge, and so there is exactly ONE jail implementation.
    """
    stem = _slug(name)
    if not stem.lower().endswith(ext):
        stem = f"{stem}{ext}"
    base = folder or default_folder()
    # must_exist=False: we are creating the file.
    path, err = resolver(os.path.join(base, stem), must_exist=False)
    if err:
        return None, err
    if not path.lower().endswith(ext):
        return None, f"Refused: output must end in {ext}."
    parent = os.path.dirname(path)
    try:
        os.makedirs(parent, exist_ok=True)
    except Exception as e:
        return None, f"Could not create {parent}: {e}"
    # Never silently overwrite: add -2, -3 ... like a browser download.
    if os.path.exists(path):
        stem2, e2 = os.path.splitext(path)
        n = 2
        while os.path.exists(f"{stem2}-{n}{e2}") and n < 200:
            n += 1
        path = f"{stem2}-{n}{e2}"
    return path, None


def default_folder():
    """
    Where generated files go by default.

    The user asked for a dedicated AURA folder rather than dumping files on the
    Desktop, and for the destination to be configurable. This is only the
    default; the UI passes an explicit `folder` when the user sets one.
    """
    home = os.path.expanduser("~")
    downloads = os.path.join(home, "Downloads")
    base = downloads if os.path.isdir(downloads) else home
    return os.path.join(base, "AURA")


def _clip(v, n=MAX_TEXT):
    return str(v if v is not None else "")[:n]


# ══════════════════════════════════════════════════════════════════════
#  PPTX — professional renderer
# ══════════════════════════════════════════════════════════════════════
#
# Every slide kind gets an intentionally-designed layout (spec §11), drawn on
# blank canvases for full control: hero title, section divider, bullets,
# two-column, process, timeline, stats, comparison table, quote, conclusion,
# references. One theme system keeps typography/color coherent, and after
# saving we RE-OPEN the file and validate it (spec §16) — "written to disk"
# is not the same as "a good deck".

_SLIDE_W = 13.333
_SLIDE_H = 7.5

# Theme palettes: (bg, ink, dim, accent, panel). The dark one is AURA's
# identity (near-black with Command Gold, like the UI theme).
THEMES = {
    "professional-dark": {
        "bg": (0x0B, 0x10, 0x1A), "ink": (0xF2, 0xF5, 0xF9),
        "dim": (0xA8, 0xB4, 0xC4), "accent": (0xE8, 0xB7, 0x4A),
        "panel": (0x14, 0x1C, 0x29),
    },
    "professional-light": {
        "bg": (0xFF, 0xFF, 0xFF), "ink": (0x10, 0x18, 0x24),
        "dim": (0x5A, 0x6B, 0x80), "accent": (0xB8, 0x86, 0x1F),
        "panel": (0xF2, 0xF4, 0xF7),
    },
    "academic": {
        "bg": (0xFA, 0xF9, 0xF6), "ink": (0x1B, 0x2A, 0x4A),
        "dim": (0x5C, 0x66, 0x74), "accent": (0x1F, 0x4E, 0x9C),
        "panel": (0xEE, 0xEC, 0xE4),
    },
    "minimal": {
        "bg": (0xFF, 0xFF, 0xFF), "ink": (0x16, 0x16, 0x16),
        "dim": (0x77, 0x77, 0x77), "accent": (0x44, 0x44, 0x44),
        "panel": (0xF4, 0xF4, 0xF4),
    },
    # Festive — "holiday homework" decks get seasonal gold-on-deep-green.
    "holiday": {
        "bg": (0x10, 0x2B, 0x1C), "ink": (0xFF, 0xF7, 0xE6),
        "dim": (0xBC, 0xD3, 0xBD), "accent": (0xF2, 0xB1, 0x3D),
        "panel": (0x1B, 0x3E, 0x2C),
    },
    # AURA identity — near-black with Command Gold + cyan for neon accents.
    "neon": {
        "bg": (0x07, 0x0A, 0x12), "ink": (0xED, 0xF3, 0xFF),
        "dim": (0x7D, 0x8D, 0xA3), "accent": (0x38, 0xBD, 0xF8),
        "panel": (0x0F, 0x14, 0x20),
    },
}
FONT_HEAD = "Calibri Light"
FONT_BODY = "Calibri"


def _theme(spec):
    name = str(spec.get("theme") or "").lower()
    for key in THEMES:
        if key in name:
            return THEMES[key], key
    if "light" in name:
        return THEMES["professional-light"], "professional-light"
    return THEMES["professional-dark"], "professional-dark"


def _rgb(t):
    return RGBColor(*t)


def _add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(fill)
    if line:
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _add_text(slide, x, y, w, h, text, size=18, color=(0, 0, 0), bold=False,
              align=None, font=FONT_BODY, italic=False, spacing=1.0):
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = _clip(text, MAX_TEXT)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = _rgb(color)
    return tb


def _add_bullets(slide, x, y, w, h, bullets, size, color, accent, font=FONT_BODY):
    """Bulleted body with the accent dot typography keeps consistent."""
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in list(bullets or [])[:9]:
        text = _clip(b, 400).strip()
        if not text:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.line_spacing = 1.08
        p.space_after = Pt(6)
        rd = p.add_run()
        rd.text = "▪  "
        rd.font.size = Pt(size)
        rd.font.name = font
        rd.font.color.rgb = _rgb(accent)
        rr = p.add_run()
        rr.text = text
        rr.font.size = Pt(size)
        rr.font.name = font
        rr.font.color.rgb = _rgb(color)
        first = False
    return tb


def _slide_bg(slide, bg):
    _add_rect(slide, 0, 0, _SLIDE_W, _SLIDE_H, bg)


def _header(slide, title, t, number=None, total=None, deck=""):
    """Title + accent rule + footer; used by all content-kind layouts."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    _add_text(slide, 0.55, 0.32, 11.5, 0.9, title, size=27, color=t["ink"],
              bold=True, font=FONT_HEAD)
    _add_rect(slide, 0.6, 1.12, 1.6, 0.045, t["accent"])
    if number and total:
        _add_text(slide, 12.1, 7.02, 1.0, 0.35, f"{number} / {total}", size=10,
                  color=t["dim"], align=PP_ALIGN.RIGHT)
    if deck:
        _add_text(slide, 0.55, 7.02, 6.0, 0.35, _clip(deck, 60), size=10, color=t["dim"])


def _notes(slide, raw):
    notes = raw.get("notes")
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = _clip(notes, 2000)
        except Exception:
            pass


def _content_fallback_bullets(raw):
    """Old/loose specs: mine ANY content field into bullets so no slide is empty."""
    bullets = list(raw.get("bullets") or [])
    for st in (raw.get("steps") or []):
        bullets.append(str(st))
    for tl in (raw.get("timeline") or []):
        lab, txt = (tl.get("label") or ""), (tl.get("text") or "")
        bullets.append(f"{lab} — {txt}".strip(" —"))
    for st in (raw.get("stats") or []):
        val, lab = (st.get("value") or ""), (st.get("label") or "")
        bullets.append(f"{val}: {lab}".strip(": "))
    col = raw.get("columns") or {}
    for side in ("left", "right"):
        c = col.get(side) or {}
        if c.get("title"):
            bullets.append(str(c["title"]))
        bullets.extend(str(b) for b in (c.get("bullets") or []))
    if raw.get("quote"):
        bullets.append(str(raw["quote"]))
    return [_clip(b, 400) for b in bullets if str(b).strip()][:9]


def _render_hero(prs, title, subtitle, t, raw=None):
    from pptx.enum.text import PP_ALIGN
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _add_rect(sl, 0, 2.62, 0.18, 2.2, t["accent"])
    kicker = (raw or {}).get("purpose") or "A presentation by NOVA"
    _add_text(sl, 0.9, 2.15, 11.5, 0.5, str(kicker).upper(), size=13,
              color=t["accent"], bold=True)
    _add_text(sl, 0.9, 2.62, 11.8, 1.9, title, size=52, color=t["ink"],
              bold=True, font=FONT_HEAD)
    if subtitle:
        _add_text(sl, 0.9, 4.45, 11.2, 0.9, subtitle, size=18, color=t["dim"], spacing=1.1)
    _add_text(sl, 0.9, 6.6, 8.0, 0.4, datetime.date.today().strftime("%d %B %Y"),
              size=12, color=t["dim"], align=PP_ALIGN.LEFT)
    if raw:
        _notes(sl, raw)
    return sl


def _render_section(prs, raw, idx, total, deck, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["panel"])
    _add_rect(sl, 0, 0, 0.35, _SLIDE_H, t["accent"])
    _add_text(sl, 1.1, 2.4, 2.5, 1.6, f"{idx:02d}", size=80, color=t["accent"],
              bold=True, font=FONT_HEAD)
    _add_text(sl, 3.3, 2.68, 9.3, 1.4, raw.get("title") or f"Section {idx}",
              size=40, color=t["ink"], bold=True, font=FONT_HEAD)
    if raw.get("purpose"):
        _add_text(sl, 3.35, 4.0, 8.8, 0.8, raw["purpose"], size=16, color=t["dim"], spacing=1.1)
    bl = raw.get("bullets") or []
    if bl:
        _add_bullets(sl, 3.35, 4.75, 8.8, 2.0, bl[:4], 15, t["ink"], t["accent"])
    _notes(sl, raw)
    return sl


def _load_image(src, resolver=None):
    """
    Resolve an image source to a real file on disk BEFORE rendering.
    http(s) → downloaded to a temp file; local path → resolver (home jail).
    Returns (path_or_None, ok, note).
    """
    s = str(src or "").strip()
    if not s or len(s) > 400:
        return None, False, "empty image source"
    if s.startswith(("http://", "https://")):
        import urllib.request
        import tempfile
        try:
            req = urllib.request.Request(s, headers={"User-Agent": "NOVA/1.0"})
            with urllib.request.urlopen(req, timeout=15) as r:
                data = r.read(25 * 1024 * 1024)
        except Exception as e:
            return None, False, f"download failed: {e}"
        ext = (s.split("?")[0].split(".")[-1] or "jpg").lower()
        if ext not in ("jpg", "jpeg", "png", "gif", "webp", "bmp"):
            ext = "jpg"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix="." + ext)
        tmp.write(data)
        tmp.close()
        return tmp.name, True, "downloaded"
    if resolver is None:
        return None, False, "local image path but no path resolver"
    try:
        resolved = resolver(s)
        # bridge._resolve_path returns (path, err); callers may also give a
        # plain path-returning callable.
        if isinstance(resolved, tuple):
            path, err = resolved[0], resolved[1]
            if err:
                return None, False, err
            if not path:
                return None, False, "path refused by the resolver"
            return str(path), True, "resolved"
        return str(resolved), True, "resolved"
    except Exception as e:
        return None, False, str(e)


def _image_px(path):
    """
    Read PNG/JPEG/GIF header dimensions in pure Python (no PIL dependency at
    runtime; PIL is a test-time cross-check only). Returns (w, h) or (None, None).
    """
    try:
        with open(path, "rb") as fh:
            head = fh.read(64)
    except Exception:
        return None, None
    if not head:
        return None, None
    try:
        if head[:8] == b"\x89PNG\r\n\x1a\n" and head[12:16] == b"IHDR":
            import struct as _st
            w, h = _st.unpack(">II", head[16:24])
            return int(w), int(h)
        if head[:3] == b"\xff\xd8\xff":
            # Scan JPEG segments for SOFn.
            import struct as _st
            i, n = 2, len(head)
            while i < n:
                if head[i] != 0xFF or i + 9 > n:
                    i += 1; continue
                marker = head[i + 1]
                if marker in (0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7,
                              0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF):
                    h, w = _st.unpack(">HH", head[i + 5:i + 9])
                    return int(w), int(h)
                seg = _st.unpack(">H", head[i + 2:i + 4])[0]
                i += 2 + seg
        if head[:6] in (b"GIF87a", b"GIF89a"):
            import struct as _st
            w, h = _st.unpack("<HH", head[6:10])
            return int(w), int(h)
    except Exception:
        pass
    return None, None


def _render_image(prs, raw, idx, total, deck, t):
    """'image' slides: full-bleed-ish picture + caption. No image → honest
    placeholder text instead of silently dropping the slide."""
    from pptx.util import Inches
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or f"Slide {idx}", t, idx, total, deck)
    path = raw.get("_imagePath")
    note = raw.get("_imageNote") or "no image source supplied"
    if not path or not os.path.exists(path):
        _add_bullets(sl, 0.9, 1.9, 11.4, 4.4,
                     [f"Image unavailable — {note}.",
                      "Add an image (path or https URL) to this slide in the outline to embed it."],
                     18, t["ink"], t["accent"])
        _notes(sl, raw)
        return
    # Fit inside the content box, preserving aspect ratio. Dimensions come
    # from our own header parser (no pptx.image module exists in python-pptx
    # 0.6.x — importing it was the crash that made the whole deck unsavable).
    pw, ph = _image_px(path)
    box_w, box_h = _SLIDE_W - 1.8, 4.4
    if pw and ph:
        ar = pw / ph
        if ar >= box_w / box_h:
            w_i, h_i = box_w, box_w / ar
        else:
            h_i, w_i = box_h, box_h * ar
        left = (_SLIDE_W - w_i) / 2
        try:
            sl.shapes.add_picture(path, Inches(left), Inches(1.62),
                                  width=Inches(w_i), height=Inches(h_i))
        except Exception as e:
            _add_bullets(sl, 0.9, 1.9, 11.4, 4.4,
                         [f"Image could not be rendered — {e}."], 18, t["ink"], t["accent"])
            _notes(sl, raw)
            return
    else:
        # Unknown format: height-only keeps the aspect ratio (no distortion).
        sl.shapes.add_picture(path, Inches(0.9), Inches(1.62),
                              height=Inches(box_h))
        _add_bullets(sl, 0.9, 1.9, 11.4, 4.4,
                     [f"Image could not be rendered — {e}."], 18, t["ink"], t["accent"])
        _notes(sl, raw)
        return
    cap = raw.get("imageCaption") or raw.get("purpose") or ""
    if cap:
        _add_text(sl, 0.9, 6.3, _SLIDE_W - 1.8, 0.7, str(cap), size=13,
                  color=t["dim"], align=None)
    _notes(sl, raw)
    return sl


def _render_bullets(prs, raw, idx, total, deck, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or f"Slide {idx}", t, idx, total, deck)
    bullets = _content_fallback_bullets(raw)
    _add_bullets(sl, 0.75, 1.55, 11.9, 5.2, bullets, 18, t["ink"], t["accent"])
    _notes(sl, raw)
    return sl


def _render_two_column(prs, raw, idx, total, deck, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or f"Slide {idx}", t, idx, total, deck)
    col = raw.get("columns") or {}
    left, right = col.get("left") or {}, col.get("right") or {}
    lb = list(left.get("bullets") or []) or [b for b in (raw.get("bullets") or [])][:4]
    rb = list(right.get("bullets") or []) or [b for b in (raw.get("bullets") or [])][4:8]
    # Left card
    _add_rect(sl, 0.6, 1.5, 5.95, 5.1, t["panel"])
    _add_text(sl, 0.9, 1.7, 5.4, 0.5, left.get("title") or "", size=16,
              color=t["accent"], bold=True)
    _add_bullets(sl, 0.9, 2.3, 5.35, 4.1, lb[:6], 15, t["ink"], t["accent"])
    # Right card
    _add_rect(sl, 6.85, 1.5, 5.95, 5.1, t["panel"])
    _add_text(sl, 7.15, 1.7, 5.4, 0.5, right.get("title") or "", size=16,
              color=t["accent"], bold=True)
    _add_bullets(sl, 7.15, 2.3, 5.35, 4.1, rb[:6], 15, t["ink"], t["accent"])
    _notes(sl, raw)
    return sl


def _render_process(prs, raw, idx, total, deck, t):
    from pptx.enum.text import PP_ALIGN
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or f"Slide {idx}", t, idx, total, deck)
    steps = [str(s) for s in (raw.get("steps") or raw.get("bullets") or []) if str(s).strip()][:6]
    n = max(1, len(steps))
    gap, w = 0.25, min(2.2, (12.2 - 0.25 * (n - 1)) / n)
    x = 0.6
    for i, s in enumerate(steps):
        _add_rect(sl, x, 2.0, w, 0.75, t["accent"])
        _add_text(sl, x, 2.08, w, 0.55, f"Step {i + 1}", size=14, color=t["bg"],
                  bold=True, align=PP_ALIGN.CENTER)
        _add_rect(sl, x, 2.75, w, 2.9, t["panel"])
        _add_text(sl, x + 0.12, 2.95, w - 0.24, 2.5, _clip(s, 220), size=13,
                  color=t["ink"], align=PP_ALIGN.LEFT, spacing=1.05)
        x += w + gap
    _notes(sl, raw)
    return sl


def _render_timeline(prs, raw, idx, total, deck, t):
    from pptx.enum.text import PP_ALIGN
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or f"Slide {idx}", t, idx, total, deck)
    items = [x for x in (raw.get("timeline") or []) if (x.get("label") or x.get("text"))][:6]
    if items:
        y_line = 3.55
        _add_rect(sl, 0.8, y_line, 11.7, 0.05, t["accent"])
        n = len(items)
        step = 11.7 / max(1, n)
        for i, it in enumerate(items):
            cx = 0.8 + step * i + step / 2
            _add_rect(sl, cx - 0.07, y_line - 0.055, 0.16, 0.16, t["accent"])
            above = (i % 2 == 0)
            _add_text(sl, cx - step / 2 + 0.1, (y_line - 1.2) if above else (y_line + 0.35),
                      step - 0.2, 0.5, _clip(it.get("label") or "", 60), size=14,
                      color=t["accent"], bold=True, align=PP_ALIGN.CENTER)
            _add_text(sl, cx - step / 2 + 0.1, (y_line - 0.72) if above else (y_line + 0.85),
                      step - 0.2, 1.5, _clip(it.get("text") or "", 200), size=11.5,
                      color=t["ink"], align=PP_ALIGN.CENTER, spacing=1.02)
    _notes(sl, raw)
    return sl


def _render_stats(prs, raw, idx, total, deck, t):
    from pptx.enum.text import PP_ALIGN
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or f"Slide {idx}", t, idx, total, deck)
    stats = [s for s in (raw.get("stats") or []) if (s.get("value") or s.get("label"))][:4]
    n = max(1, len(stats))
    w = min(2.9, (12.2 - 0.3 * (n - 1)) / n)
    x = 0.6
    for s in stats:
        _add_rect(sl, x, 1.9, w, 3.3, t["panel"])
        _add_rect(sl, x, 1.9, w, 0.07, t["accent"])
        _add_text(sl, x + 0.1, 2.25, w - 0.2, 1.2, _clip(s.get("value") or "", 24),
                  size=40, color=t["accent"], bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
        _add_text(sl, x + 0.18, 3.55, w - 0.36, 1.5, _clip(s.get("label") or "", 200),
                  size=13, color=t["ink"], align=PP_ALIGN.CENTER, spacing=1.05)
        x += w + 0.3
    if raw.get("bullets"):
        _add_bullets(sl, 0.75, 5.45, 11.9, 1.3, raw["bullets"][:3], 14, t["ink"], t["accent"])
    _notes(sl, raw)
    return sl


def _render_comparison(prs, raw, idx, total, deck, t):
    table = raw.get("table") or {}
    cols = [str(c) for c in (table.get("columns") or [])][:8]
    rows = table.get("rows") or []
    if not cols or not rows:
        return _render_two_column(prs, raw, idx, total, deck, t)
    from pptx.util import Inches, Pt
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or f"Slide {idx}", t, idx, total, deck)
    rr = min(len(rows), 8)
    shp = sl.shapes.add_table(rr + 1, len(cols), Inches(0.6), Inches(1.55),
                              Inches(12.1), Inches(min(4.9, 0.55 * (rr + 1))))
    tbl = shp.table
    for j, c in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = _clip(c, 80)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14); r.font.bold = True
                r.font.color.rgb = _rgb(t["bg"] if t["bg"] != (0xFF, 0xFF, 0xFF) else t["ink"])
    for i in range(rr):
        for j in range(len(cols)):
            cell = tbl.cell(i + 1, j)
            val = rows[i][j] if j < len(rows[i]) else ""
            cell.text = _clip(val, 200)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12.5)
                    r.font.color.rgb = _rgb(t["ink"])
    _notes(sl, raw)
    return sl


def _render_quote(prs, raw, idx, total, deck, t):
    from pptx.enum.text import PP_ALIGN
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _add_rect(sl, 0.9, 2.6, 0.16, 2.4, t["accent"])
    quote = raw.get("quote") or " ".join(_content_fallback_bullets(raw)[:2])
    _add_text(sl, 1.35, 2.5, 10.6, 2.6, f"“{_clip(quote, 400)}”", size=28,
              color=t["ink"], italic=True, font=FONT_HEAD, align=PP_ALIGN.LEFT, spacing=1.15)
    if raw.get("attribution"):
        _add_text(sl, 1.4, 5.2, 9.0, 0.5, "— " + _clip(raw["attribution"], 120),
                  size=15, color=t["dim"])
    if raw.get("title") and raw["title"] != "Quote":
        _add_text(sl, 0.9, 1.5, 11.5, 0.6, raw["title"], size=15, color=t["accent"], bold=True)
    _notes(sl, raw)
    return sl


def _render_conclusion(prs, raw, idx, total, deck, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _add_rect(sl, 0, 0, _SLIDE_W, 0.16, t["accent"])
    _header(sl, raw.get("title") or "Conclusion", t, idx, total, deck)
    bullets = _content_fallback_bullets(raw)
    _add_bullets(sl, 0.9, 1.75, 11.5, 4.6, bullets, 19, t["ink"], t["accent"])
    _notes(sl, raw)
    return sl


def _render_references(prs, raw, idx, total, deck, t):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, t["bg"])
    _header(sl, raw.get("title") or "References", t, idx, total, deck)
    refs = _content_fallback_bullets(raw)
    _add_bullets(sl, 0.75, 1.6, 11.9, 5.0, refs, 13.5, t["dim"], t["accent"])
    _notes(sl, raw)
    return sl


_KIND_RENDERERS = {
    "section": _render_section,
    "bullets": _render_bullets,
    "image": _render_image,
    "two-column": _render_two_column,
    "process": _render_process,
    "timeline": _render_timeline,
    "stats": _render_stats,
    "comparison": _render_comparison,
    "quote": _render_quote,
    "conclusion": _render_conclusion,
    "references": _render_references,
}


def validate_pptx(path, expected_slides=None):
    """
    Post-write artifact validation (spec §16): re-open the FILE and check the
    things a user would notice — slide count, missing titles, slides whose
    body is empty, giant over-stuffed bodies. Returns a report, never throws.
    """
    report = {"ok": False, "fileExists": os.path.isfile(path), "slideCount": 0,
              "issues": []}
    if not report["fileExists"]:
        report["issues"].append("file not found after save")
        return report
    try:
        prs = Presentation(path)
    except Exception as e:
        report["issues"].append(f"file cannot be re-opened: {e}")
        return report
    slides = list(prs.slides)
    report["slideCount"] = len(slides)
    if expected_slides is not None and len(slides) != expected_slides:
        report["issues"].append(f"slide count {len(slides)} != expected {expected_slides}")
    if not slides:
        report["issues"].append("file opens but has no slides")
        return report
    for i, s in enumerate(slides):
        texts = [sh.text_frame.text.strip() for sh in s.shapes
                 if getattr(sh, "has_text_frame", False)]
        joined = " ".join(x for x in texts if x).strip()
        if not joined:
            report["issues"].append(f"slide {i + 1} has no text at all")
        elif i > 0 and len(joined) < 8:
            report["issues"].append(f"slide {i + 1} is nearly empty ({len(joined)} chars)")
        if sum(len(x) for x in texts) > 2600:
            report["issues"].append(f"slide {i + 1} likely overflows ({sum(len(x) for x in texts)} chars)")
    report["ok"] = not report["issues"]
    return report


def build_pptx(spec, folder=None, resolver=None):
    """
    spec = {
      "title": str, "subtitle": str, "theme": str (optional),
      "slides": [ {"kind": title|section|bullets|two-column|process|timeline|
                          stats|comparison|quote|conclusion|references,
                   "title": str, "purpose": str, "bullets": [str], "columns": …,
                   "steps": …, "timeline": …, "stats": …, "table": …,
                   "quote": str, "notes": str}, ... ]
    }
    Old specs (no kinds) render exactly as before: hero + bullet slides.
    """
    if not HAS_PPTX:
        return {"ok": False, "missing": "python-pptx",
                "message": "PowerPoint generation needs python-pptx.  "
                           "pip install python-pptx"}
    if not isinstance(spec, dict):
        return {"ok": False, "message": "Bad spec: expected an object."}

    slides = spec.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return {"ok": False, "message": "The outline has no slides."}
    if len(slides) > MAX_SLIDES:
        return {"ok": False,
                "message": f"Refused: {len(slides)} slides exceeds the {MAX_SLIDES} cap."}
    # Drop slides that carry nothing at all (no title AND no content).
    usable = [s for s in slides
              if isinstance(s, dict) and (str(s.get("title") or "").strip()
                                          or _content_fallback_bullets(s)
                                          or (s.get("table") or {}).get("rows"))]
    if not usable:
        return {"ok": False, "message": "The outline has no usable slides."}

    title = _clip(spec.get("title") or "Untitled", 160)
    path, err = _target_path(folder, spec.get("filename") or title, ".pptx", resolver)
    if err:
        return {"ok": False, "message": err}

    t, theme_name = _theme(spec)
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)

    first_is_title = str(usable[0].get("kind") or "").lower() == "title"
    deck_title = title
    made = 0

    # Hero: from the spec's title slide, or synthesised from spec.title so
    # legacy outlines (title + N bullet slides) keep their old N+1 shape.
    if first_is_title:
        raw0 = usable[0]
        _render_hero(prs, _clip(raw0.get("title") or title, 160),
                     _clip(raw0.get("subtitle") or spec.get("subtitle") or "", 200),
                     t, raw0)
        content = usable[1:]
        made = 1
    else:
        _render_hero(prs, title,
                     _clip(spec.get("subtitle") or datetime.date.today().strftime("%d %B %Y"), 200),
                     t, None)
        content = usable
        made = 1

    total = made + len(content)
    rendered = 0
    embedded, failed = [], []
    for offset, raw in enumerate(content):
        idx = made + offset + 1
        kind = str(raw.get("kind") or "").lower().strip()
        renderer = _KIND_RENDERERS.get(kind, _render_bullets)
        render_raw = raw
        if raw.get("image"):
            # NOTE: img_* names — `path` is the OUTPUT deck path, shadowing it
            # here made prs.save() write the deck over the image file.
            img_path, img_ok, img_note = _load_image(raw["image"], resolver)
            render_raw = {**raw, "_imagePath": img_path if img_ok else None,
                          "_imageNote": img_note}
            if img_ok:
                embedded.append(f"slide {idx}")
            else:
                failed.append(f"slide {idx}: {img_note}")
        try:
            renderer(prs, render_raw, idx, total, deck_title, t)
        except Exception:
            # A fancy layout failing must never kill the deck: render plain.
            fb = prs.slides.add_slide(prs.slide_layouts[6])
            _slide_bg(fb, t["bg"])
            _header(fb, raw.get("title") or f"Slide {idx}", t, idx, total, deck_title)
            _add_bullets(fb, 0.75, 1.55, 11.9, 5.2, _content_fallback_bullets(raw),
                         18, t["ink"], t["accent"])
            _notes(fb, raw)
        rendered += 1

    try:
        prs.save(path)
    except Exception as e:
        return {"ok": False, "message": f"Could not save: {e}"}

    # Artifact verification — report honestly instead of blind success (§16).
    validation = validate_pptx(path, expected_slides=made + rendered)
    if failed:
        issues = list(validation.get("issues") or [])
        issues.insert(0, f"{len(failed)} image(s) could not be embedded: {'; '.join(failed[:3])}")
        validation["issues"] = issues
    validation["embedded_images"] = embedded
    validation["failed_images"] = failed
    total_slides = made + rendered
    msg = f"Created a {total_slides}-slide presentation ({theme_name}) at {path}"
    if embedded:
        msg += f" — {len(embedded)} image(s) embedded."
    if validation["ok"]:
        msg += " — validated: every slide has content."
    else:
        msg += f" — validation warnings: {'; '.join(validation['issues'][:3])}"

    return {"ok": True, "path": path, "kind": "pptx",
            "slides": total_slides, "bytes": os.path.getsize(path),
            "theme": theme_name, "validation": validation,
            "embedded_images": len(embedded), "failed_images": failed,
            "message": msg}


# ══════════════════════════════════════════════════════════════════════
#  XLSX
# ══════════════════════════════════════════════════════════════════════

def build_xlsx(spec, folder=None, resolver=None):
    """
    spec = {
      "title": str,
      "sheets": [ {"name": str, "columns": [str], "rows": [[cell, ...]]} ]
    }
    A single-sheet shorthand ({"columns":…, "rows":…}) is also accepted.
    """
    if not HAS_XLSX:
        return {"ok": False, "missing": "openpyxl",
                "message": "Spreadsheet generation needs openpyxl.  pip install openpyxl"}
    if not isinstance(spec, dict):
        return {"ok": False, "message": "Bad spec: expected an object."}

    sheets = spec.get("sheets")
    if not sheets and (spec.get("rows") or spec.get("columns")):
        sheets = [{"name": spec.get("title") or "Sheet1",
                   "columns": spec.get("columns"), "rows": spec.get("rows")}]
    if not isinstance(sheets, list) or not sheets:
        return {"ok": False, "message": "The outline has no sheets."}

    title = _clip(spec.get("title") or "Untitled", 160)
    path, err = _target_path(folder, spec.get("filename") or title, ".xlsx", resolver)
    if err:
        return {"ok": False, "message": err}

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="1F3B52")
    total_rows = 0

    for i, sh in enumerate(sheets[:12]):
        if not isinstance(sh, dict):
            continue
        # Excel sheet names: 31 chars, and []:*?/\ are illegal.
        name = re.sub(r"[\[\]:*?/\\]", "-", str(sh.get("name") or f"Sheet{i+1}"))[:31] or f"Sheet{i+1}"
        ws = wb.create_sheet(name)
        cols = sh.get("columns") or []
        rows = sh.get("rows") or []
        if len(rows) > MAX_ROWS:
            return {"ok": False,
                    "message": f"Refused: {len(rows)} rows exceeds the {MAX_ROWS} cap."}

        r = 1
        if cols:
            for c, label in enumerate(list(cols)[:MAX_COLS], start=1):
                cell = ws.cell(row=1, column=c, value=_clip(label, 200))
                cell.font = head_font
                cell.fill = head_fill
                cell.alignment = Alignment(horizontal="center")
            ws.freeze_panes = "A2"
            r = 2

        for row in rows:
            if not isinstance(row, (list, tuple)):
                row = [row]
            for c, val in enumerate(list(row)[:MAX_COLS], start=1):
                # Numbers stay numbers so Excel can sum them; everything else
                # becomes text. A string like "=cmd" would be a formula
                # injection, so any leading = + - @ is prefixed with '.
                if isinstance(val, bool) or val is None:
                    out = "" if val is None else str(val)
                elif isinstance(val, (int, float)):
                    out = val
                else:
                    out = _clip(val, 2000)
                    if out[:1] in ("=", "+", "-", "@"):
                        out = "'" + out
                ws.cell(row=r, column=c, value=out)
            r += 1
            total_rows += 1

        # Readable widths, capped so one long cell cannot make a 400px column.
        widths = {}
        for row in ws.iter_rows(min_row=1, max_row=min(r, 200)):
            for cell in row:
                if cell.value is None:
                    continue
                widths[cell.column] = min(48, max(widths.get(cell.column, 10),
                                                  len(str(cell.value)) + 2))
        for col, w in widths.items():
            ws.column_dimensions[get_column_letter(col)].width = w

    if not wb.sheetnames:
        return {"ok": False, "message": "Nothing to write — every sheet was empty."}

    try:
        wb.save(path)
    except Exception as e:
        return {"ok": False, "message": f"Could not save: {e}"}

    return {"ok": True, "path": path, "kind": "xlsx",
            "sheets": len(wb.sheetnames), "rows": total_rows,
            "bytes": os.path.getsize(path),
            "message": f"Created a {len(wb.sheetnames)}-sheet workbook "
                       f"({total_rows} rows) at {path}"}


# ══════════════════════════════════════════════════════════════════════
#  DOCX
# ══════════════════════════════════════════════════════════════════════

def build_docx(spec, folder=None, resolver=None):
    """
    spec = {
      "title": str, "subtitle": str,
      "sections": [ {"heading": str, "level": 1..3,
                     "paragraphs": [str], "bullets": [str]} ]
    }
    """
    if not HAS_DOCX:
        return {"ok": False, "missing": "python-docx",
                "message": "Document generation needs python-docx.  pip install python-docx"}
    if not isinstance(spec, dict):
        return {"ok": False, "message": "Bad spec: expected an object."}

    sections = spec.get("sections") or []
    if isinstance(sections, dict):
        sections = [sections]
    if not sections:
        return {"ok": False, "message": "The outline has no sections."}
    if len(sections) > MAX_SECTIONS:
        return {"ok": False,
                "message": f"Refused: {len(sections)} sections exceeds the {MAX_SECTIONS} cap."}

    title = _clip(spec.get("title") or "Untitled", 160)
    path, err = _target_path(folder, spec.get("filename") or title, ".docx", resolver)
    if err:
        return {"ok": False, "message": err}

    d = docx.Document()
    h = d.add_heading(title, level=0)
    try:
        h.runs[0].font.color.rgb = DocxRGB(*ACCENT)
    except Exception:
        pass
    if spec.get("subtitle"):
        p = d.add_paragraph(_clip(spec["subtitle"], 400))
        try:
            p.runs[0].font.color.rgb = DocxRGB(*DIM)
            p.runs[0].font.size = DocxPt(11)
        except Exception:
            pass

    words = 0
    for sec in sections:
        if not isinstance(sec, dict):
            continue
        if sec.get("heading"):
            lvl = sec.get("level")
            lvl = lvl if isinstance(lvl, int) and 1 <= lvl <= 4 else 1
            d.add_heading(_clip(sec["heading"], 200), level=lvl)
        paras = sec.get("paragraphs") or []
        if isinstance(paras, str):
            paras = [paras]
        for para in list(paras)[:60]:
            text = _clip(para, MAX_TEXT)
            if text.strip():
                d.add_paragraph(text)
                words += len(text.split())
        bullets = sec.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for b in list(bullets)[:60]:
            text = _clip(b, 1000)
            if text.strip():
                d.add_paragraph(text, style="List Bullet")
                words += len(text.split())

    try:
        d.save(path)
    except Exception as e:
        return {"ok": False, "message": f"Could not save: {e}"}

    return {"ok": True, "path": path, "kind": "docx",
            "sections": len(sections), "words": words,
            "bytes": os.path.getsize(path),
            "message": f"Created a {len(sections)}-section document "
                       f"(~{words} words) at {path}"}


BUILDERS = {"pptx": build_pptx, "xlsx": build_xlsx, "docx": build_docx}


def build(kind, spec, folder=None, resolver=None):
    """Dispatch to a builder. Unknown kinds are refused, not guessed."""
    fn = BUILDERS.get(str(kind or "").lower().lstrip("."))
    if not fn:
        return {"ok": False,
                "message": f"Unknown document type '{kind}'. "
                           f"Supported: {', '.join(sorted(BUILDERS))}."}
    if resolver is None:
        return {"ok": False, "message": "Internal: no path resolver supplied."}
    return fn(spec, folder=folder, resolver=resolver)
