"""
AURA :: Document Builder
========================
Turns a structured outline into a REAL Office file: .pptx, .xlsx or .docx.

WHAT THIS MODULE IS, AND IS NOT
-------------------------------
It is the *rendering* half only. It takes an already-structured outline —
a dict of slides / rows / sections — and writes a genuine OOXML file to disk.

It deliberately contains **no AI**. The prompt-to-outline step happens in the
browser (`js/ai/doc-agent.js`) using whichever provider the user already
configured, and its JSON is validated here before a single byte is written.
That split matters: the dangerous part (writing files) stays small, auditable
and fully testable without a model running.

HONESTY
-------
python-pptx / openpyxl / python-docx are OPTIONAL. If one is missing, the
matching builder reports exactly which `pip install` fixes it and refuses.
It never writes a fake file, never writes a .txt and calls it a .pptx.

SAFETY
------
Every output path goes through `bridge._resolve_path`, so the same jail that
protects the file plugin protects this one: no writing outside the user's
folders, no symlink escape, no credential directories. The extension is forced
to match the builder, so a "presentation" can never be written as `.exe`.
"""

import os
import re
import datetime

# The builders are optional; report per-format rather than failing at import.
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except Exception:                                  # pragma: no cover
    HAS_PPTX = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter
    HAS_XLSX = True
except Exception:                                  # pragma: no cover
    HAS_XLSX = False

try:
    import docx
    from docx.shared import Pt as DocxPt, RGBColor as DocxRGB
    HAS_DOCX = True
except Exception:                                  # pragma: no cover
    HAS_DOCX = False


MAX_SLIDES = 40
MAX_ROWS = 5000
MAX_COLS = 60
MAX_SECTIONS = 80
MAX_TEXT = 4000

# AURA's accent, so generated decks look like they came from AURA.
ACCENT = (0x38, 0xBD, 0xF8)
INK = (0x10, 0x18, 0x24)
DIM = (0x5A, 0x6B, 0x80)


def capabilities():
    """What can actually be built on this machine, right now."""
    return {
        "ok": True,
        "pptx": HAS_PPTX,
        "xlsx": HAS_XLSX,
        "docx": HAS_DOCX,
        "any": HAS_PPTX or HAS_XLSX or HAS_DOCX,
        "install": {
            "pptx": "pip install python-pptx",
            "xlsx": "pip install openpyxl",
            "docx": "pip install python-docx",
        },
        "note": ("Each format needs its own library. AURA reports which one is "
                 "missing instead of writing a placeholder file."),
    }


def _slug(name, fallback="aura-document"):
    """A safe file stem. Never a path — separators are stripped, not escaped."""
    s = re.sub(r"[^\w\s.-]", "", str(name or "")).strip()
    s = re.sub(r"[\s_]+", "-", s).strip("-.")
    s = s[:60] or fallback
    # Windows reserved device names would make an unopenable file.
    if s.upper().split(".")[0] in {
            "CON", "PRN", "AUX", "NUL", "COM1", "COM2", "COM3", "COM4",
            "LPT1", "LPT2", "LPT3"}:
        s = f"aura-{s}"
    return s


def _target_path(folder, name, ext, resolver):
    """
    Build and jail-check the output path.

    `resolver` is injected (bridge._resolve_path) so this module stays testable
    without importing bridge, and so there is exactly ONE jail implementation.
    """
    stem = _slug(name)
    if not stem.lower().endswith(ext):
        stem = f"{stem}{ext}"
    base = folder or default_folder()
    # must_exist=False: we are creating the file.
    path, err = resolver(os.path.join(base, stem), must_exist=False)
    if err:
        return None, err
    if not path.lower().endswith(ext):
        return None, f"Refused: output must end in {ext}."
    parent = os.path.dirname(path)
    try:
        os.makedirs(parent, exist_ok=True)
    except Exception as e:
        return None, f"Could not create {parent}: {e}"
    # Never silently overwrite: add -2, -3 ... like a browser download.
    if os.path.exists(path):
        stem2, e2 = os.path.splitext(path)
        n = 2
        while os.path.exists(f"{stem2}-{n}{e2}") and n < 200:
            n += 1
        path = f"{stem2}-{n}{e2}"
    return path, None


def default_folder():
    """
    Where generated files go by default.

    The user asked for a dedicated AURA folder rather than dumping files on the
    Desktop, and for the destination to be configurable. This is only the
    default; the UI passes an explicit `folder` when the user sets one.
    """
    home = os.path.expanduser("~")
    downloads = os.path.join(home, "Downloads")
    base = downloads if os.path.isdir(downloads) else home
    return os.path.join(base, "AURA")


def _clip(v, n=MAX_TEXT):
    return str(v if v is not None else "")[:n]


# ══════════════════════════════════════════════════════════════════════
#  PPTX
# ══════════════════════════════════════════════════════════════════════

def build_pptx(spec, folder=None, resolver=None):
    """
    spec = {
      "title": str, "subtitle": str,
      "slides": [ {"title": str, "bullets": [str], "notes": str}, ... ]
    }
    """
    if not HAS_PPTX:
        return {"ok": False, "missing": "python-pptx",
                "message": "PowerPoint generation needs python-pptx.  "
                           "pip install python-pptx"}
    if not isinstance(spec, dict):
        return {"ok": False, "message": "Bad spec: expected an object."}

    slides = spec.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return {"ok": False, "message": "The outline has no slides."}
    if len(slides) > MAX_SLIDES:
        return {"ok": False,
                "message": f"Refused: {len(slides)} slides exceeds the {MAX_SLIDES} cap."}

    title = _clip(spec.get("title") or "Untitled", 160)
    path, err = _target_path(folder, spec.get("filename") or title, ".pptx", resolver)
    if err:
        return {"ok": False, "message": err}

    prs = Presentation()
    prs.slide_width = Inches(13.333)          # 16:9
    prs.slide_height = Inches(7.5)

    # ── title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    if len(s.placeholders) > 1:
        sub = _clip(spec.get("subtitle") or
                    datetime.date.today().strftime("%d %B %Y"), 200)
        s.placeholders[1].text = sub
    try:
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*ACCENT)
    except Exception:
        pass

    # ── content slides
    made = 0
    for raw in slides:
        if not isinstance(raw, dict):
            continue
        layout = prs.slide_layouts[1]          # Title and Content
        sl = prs.slides.add_slide(layout)
        sl.shapes.title.text = _clip(raw.get("title") or f"Slide {made + 2}", 160)

        bullets = raw.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        body = None
        for ph in sl.placeholders:
            if ph.placeholder_format.idx != 0:
                body = ph
                break
        if body is not None:
            tf = body.text_frame
            tf.clear()
            first = True
            for b in list(bullets)[:12]:
                text = _clip(b, 400)
                if not text.strip():
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(18)
                first = False
            if first:                          # no usable bullets
                tf.paragraphs[0].text = ""

        notes = raw.get("notes")
        if notes:
            try:
                sl.notes_slide.notes_text_frame.text = _clip(notes, 2000)
            except Exception:
                pass
        made += 1

    try:
        prs.save(path)
    except Exception as e:
        return {"ok": False, "message": f"Could not save: {e}"}

    return {"ok": True, "path": path, "kind": "pptx",
            "slides": made + 1, "bytes": os.path.getsize(path),
            "message": f"Created a {made + 1}-slide presentation at {path}"}


# ══════════════════════════════════════════════════════════════════════
#  XLSX
# ══════════════════════════════════════════════════════════════════════

def build_xlsx(spec, folder=None, resolver=None):
    """
    spec = {
      "title": str,
      "sheets": [ {"name": str, "columns": [str], "rows": [[cell, ...]]} ]
    }
    A single-sheet shorthand ({"columns":…, "rows":…}) is also accepted.
    """
    if not HAS_XLSX:
        return {"ok": False, "missing": "openpyxl",
                "message": "Spreadsheet generation needs openpyxl.  pip install openpyxl"}
    if not isinstance(spec, dict):
        return {"ok": False, "message": "Bad spec: expected an object."}

    sheets = spec.get("sheets")
    if not sheets and (spec.get("rows") or spec.get("columns")):
        sheets = [{"name": spec.get("title") or "Sheet1",
                   "columns": spec.get("columns"), "rows": spec.get("rows")}]
    if not isinstance(sheets, list) or not sheets:
        return {"ok": False, "message": "The outline has no sheets."}

    title = _clip(spec.get("title") or "Untitled", 160)
    path, err = _target_path(folder, spec.get("filename") or title, ".xlsx", resolver)
    if err:
        return {"ok": False, "message": err}

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="1F3B52")
    total_rows = 0

    for i, sh in enumerate(sheets[:12]):
        if not isinstance(sh, dict):
            continue
        # Excel sheet names: 31 chars, and []:*?/\ are illegal.
        name = re.sub(r"[\[\]:*?/\\]", "-", str(sh.get("name") or f"Sheet{i+1}"))[:31] or f"Sheet{i+1}"
        ws = wb.create_sheet(name)
        cols = sh.get("columns") or []
        rows = sh.get("rows") or []
        if len(rows) > MAX_ROWS:
            return {"ok": False,
                    "message": f"Refused: {len(rows)} rows exceeds the {MAX_ROWS} cap."}

        r = 1
        if cols:
            for c, label in enumerate(list(cols)[:MAX_COLS], start=1):
                cell = ws.cell(row=1, column=c, value=_clip(label, 200))
                cell.font = head_font
                cell.fill = head_fill
                cell.alignment = Alignment(horizontal="center")
            ws.freeze_panes = "A2"
            r = 2

        for row in rows:
            if not isinstance(row, (list, tuple)):
                row = [row]
            for c, val in enumerate(list(row)[:MAX_COLS], start=1):
                # Numbers stay numbers so Excel can sum them; everything else
                # becomes text. A string like "=cmd" would be a formula
                # injection, so any leading = + - @ is prefixed with '.
                if isinstance(val, bool) or val is None:
                    out = "" if val is None else str(val)
                elif isinstance(val, (int, float)):
                    out = val
                else:
                    out = _clip(val, 2000)
                    if out[:1] in ("=", "+", "-", "@"):
                        out = "'" + out
                ws.cell(row=r, column=c, value=out)
            r += 1
            total_rows += 1

        # Readable widths, capped so one long cell cannot make a 400px column.
        widths = {}
        for row in ws.iter_rows(min_row=1, max_row=min(r, 200)):
            for cell in row:
                if cell.value is None:
                    continue
                widths[cell.column] = min(48, max(widths.get(cell.column, 10),
                                                  len(str(cell.value)) + 2))
        for col, w in widths.items():
            ws.column_dimensions[get_column_letter(col)].width = w

    if not wb.sheetnames:
        return {"ok": False, "message": "Nothing to write — every sheet was empty."}

    try:
        wb.save(path)
    except Exception as e:
        return {"ok": False, "message": f"Could not save: {e}"}

    return {"ok": True, "path": path, "kind": "xlsx",
            "sheets": len(wb.sheetnames), "rows": total_rows,
            "bytes": os.path.getsize(path),
            "message": f"Created a {len(wb.sheetnames)}-sheet workbook "
                       f"({total_rows} rows) at {path}"}


# ══════════════════════════════════════════════════════════════════════
#  DOCX
# ══════════════════════════════════════════════════════════════════════

def build_docx(spec, folder=None, resolver=None):
    """
    spec = {
      "title": str, "subtitle": str,
      "sections": [ {"heading": str, "level": 1..3,
                     "paragraphs": [str], "bullets": [str]} ]
    }
    """
    if not HAS_DOCX:
        return {"ok": False, "missing": "python-docx",
                "message": "Document generation needs python-docx.  pip install python-docx"}
    if not isinstance(spec, dict):
        return {"ok": False, "message": "Bad spec: expected an object."}

    sections = spec.get("sections") or []
    if isinstance(sections, dict):
        sections = [sections]
    if not sections:
        return {"ok": False, "message": "The outline has no sections."}
    if len(sections) > MAX_SECTIONS:
        return {"ok": False,
                "message": f"Refused: {len(sections)} sections exceeds the {MAX_SECTIONS} cap."}

    title = _clip(spec.get("title") or "Untitled", 160)
    path, err = _target_path(folder, spec.get("filename") or title, ".docx", resolver)
    if err:
        return {"ok": False, "message": err}

    d = docx.Document()
    h = d.add_heading(title, level=0)
    try:
        h.runs[0].font.color.rgb = DocxRGB(*ACCENT)
    except Exception:
        pass
    if spec.get("subtitle"):
        p = d.add_paragraph(_clip(spec["subtitle"], 400))
        try:
            p.runs[0].font.color.rgb = DocxRGB(*DIM)
            p.runs[0].font.size = DocxPt(11)
        except Exception:
            pass

    words = 0
    for sec in sections:
        if not isinstance(sec, dict):
            continue
        if sec.get("heading"):
            lvl = sec.get("level")
            lvl = lvl if isinstance(lvl, int) and 1 <= lvl <= 4 else 1
            d.add_heading(_clip(sec["heading"], 200), level=lvl)
        paras = sec.get("paragraphs") or []
        if isinstance(paras, str):
            paras = [paras]
        for para in list(paras)[:60]:
            text = _clip(para, MAX_TEXT)
            if text.strip():
                d.add_paragraph(text)
                words += len(text.split())
        bullets = sec.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for b in list(bullets)[:60]:
            text = _clip(b, 1000)
            if text.strip():
                d.add_paragraph(text, style="List Bullet")
                words += len(text.split())

    try:
        d.save(path)
    except Exception as e:
        return {"ok": False, "message": f"Could not save: {e}"}

    return {"ok": True, "path": path, "kind": "docx",
            "sections": len(sections), "words": words,
            "bytes": os.path.getsize(path),
            "message": f"Created a {len(sections)}-section document "
                       f"(~{words} words) at {path}"}


BUILDERS = {"pptx": build_pptx, "xlsx": build_xlsx, "docx": build_docx}


def build(kind, spec, folder=None, resolver=None):
    """Dispatch to a builder. Unknown kinds are refused, not guessed."""
    fn = BUILDERS.get(str(kind or "").lower().lstrip("."))
    if not fn:
        return {"ok": False,
                "message": f"Unknown document type '{kind}'. "
                           f"Supported: {', '.join(sorted(BUILDERS))}."}
    if resolver is None:
        return {"ok": False, "message": "Internal: no path resolver supplied."}
    return fn(spec, folder=folder, resolver=resolver)
