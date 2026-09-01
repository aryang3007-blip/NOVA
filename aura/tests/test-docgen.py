#!/usr/bin/env python3
"""
AURA :: document builder + file organiser + QR pairing
======================================================
Proves the files are REAL Office documents (valid OOXML zips with the right
internal parts), that the path jail holds, that the organiser never moves
anything without a confirmed plan, and that the QR actually decodes.

    python3 tests/test-docgen.py
"""
import os, sys, shutil, zipfile, json, time

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
import bridge, docbuilder, organizer, devices          # noqa: E402

P, F = [], []


def ok(n, c, d=""):
    (P if c else F).append(n)
    print(("  \033[32m✓\033[0m " if c else "  \033[31m✗\033[0m ") + n
          + (f"  \033[90m{d}\033[0m" if d else ""))


def sec(t):
    print(f"\n\033[36m▸ {t}\033[0m")


R = bridge._resolve_path
OUT = os.path.join(os.path.expanduser("~"), "aura-test-out")
shutil.rmtree(OUT, ignore_errors=True)

# ════════════════════════════════════════════════════════════ CAPABILITIES
sec("CAPABILITIES ARE REPORTED HONESTLY")
caps = docbuilder.capabilities()
ok("capabilities() answers", caps["ok"])
for k in ("pptx", "xlsx", "docx"):
    ok(f"{k} availability is a real boolean", isinstance(caps[k], bool), str(caps[k]))
    ok(f"{k} has an install hint", caps["install"][k].startswith("pip install"))
ok("a default folder is offered", bool(docbuilder.default_folder()))
ok("the default folder is inside the user's own folders",
   R(docbuilder.default_folder(), must_exist=False)[1] is None)

# ════════════════════════════════════════════════════════════════ PPTX
sec("PPTX IS A REAL POWERPOINT FILE")
if caps["pptx"]:
    r = docbuilder.build("pptx", {
        "title": "AURA", "subtitle": "Generated in a test",
        "slides": [
            {"title": "Overview", "bullets": ["Runs locally", "No backend"], "notes": "hi"},
            {"title": "Numbers", "bullets": ["2365 assertions"]},
        ]}, folder=OUT, resolver=R)
    ok("build reports success", r["ok"], r.get("message", "")[:70])
    ok("the file exists on disk", os.path.isfile(r.get("path", "")), r.get("path"))
    ok("it is a valid zip (OOXML)", zipfile.is_zipfile(r["path"]))
    names = zipfile.ZipFile(r["path"]).namelist()
    ok("contains ppt/presentation.xml", "ppt/presentation.xml" in names)
    ok("contains [Content_Types].xml", "[Content_Types].xml" in names)
    ok("has 3 slides (title + 2)", r["slides"] == 3, str(r.get("slides")))
    ok("slide count matches the package",
       len([n for n in names if n.startswith("ppt/slides/slide")]) == 3,
       str([n for n in names if n.startswith("ppt/slides/slide")]))
    # The text must actually be in there, not just a blank deck.
    from pptx import Presentation
    prs = Presentation(r["path"])
    text = " ".join(sh.text_frame.text for s in prs.slides
                    for sh in s.shapes if sh.has_text_frame)
    ok("the title made it into the deck", "AURA" in text)
    ok("bullets made it into the deck", "Runs locally" in text, text[:60])
    ok("speaker notes are preserved",
       any("hi" in s.notes_slide.notes_text_frame.text for s in prs.slides
           if s.has_notes_slide))
else:
    ok("pptx refuses honestly when the library is missing",
       docbuilder.build("pptx", {"slides": [{}]}, folder=OUT, resolver=R)
       .get("missing") == "python-pptx")

# ════════════════════════════════════════════════════════════════ XLSX
sec("XLSX IS A REAL WORKBOOK")
if caps["xlsx"]:
    r = docbuilder.build("xlsx", {
        "title": "Budget",
        "sheets": [{"name": "Q1", "columns": ["Item", "Cost"],
                    "rows": [["Laptop", 900], ["Mouse", 25], ["=1+1", 3]]}],
    }, folder=OUT, resolver=R)
    ok("build reports success", r["ok"], r.get("message", "")[:70])
    ok("it is a valid zip", zipfile.is_zipfile(r["path"]))
    ok("contains xl/workbook.xml", "xl/workbook.xml" in zipfile.ZipFile(r["path"]).namelist())
    import openpyxl
    wb = openpyxl.load_workbook(r["path"])
    ws = wb["Q1"]
    ok("headers are written", ws["A1"].value == "Item" and ws["B1"].value == "Cost")
    ok("headers are bold", ws["A1"].font.bold is True)
    ok("panes are frozen below the header", ws.freeze_panes == "A2", str(ws.freeze_panes))
    ok("numbers stay numeric so Excel can sum them",
       isinstance(ws["B2"].value, (int, float)), repr(ws["B2"].value))
    ok("FORMULA INJECTION is neutralised",
       str(ws["A4"].value).startswith("'="), repr(ws["A4"].value))
else:
    ok("xlsx refuses honestly", docbuilder.build("xlsx", {"rows": [[1]]}, folder=OUT,
                                                 resolver=R).get("missing") == "openpyxl")

# ════════════════════════════════════════════════════════════════ DOCX
sec("DOCX IS A REAL WORD DOCUMENT")
if caps["docx"]:
    r = docbuilder.build("docx", {
        "title": "Report", "subtitle": "by AURA",
        "sections": [{"heading": "Intro", "level": 1,
                      "paragraphs": ["Hello world."], "bullets": ["one", "two"]}],
    }, folder=OUT, resolver=R)
    ok("build reports success", r["ok"], r.get("message", "")[:70])
    ok("it is a valid zip", zipfile.is_zipfile(r["path"]))
    ok("contains word/document.xml", "word/document.xml" in zipfile.ZipFile(r["path"]).namelist())
    import docx as _docx
    d = _docx.Document(r["path"])
    body = " ".join(p.text for p in d.paragraphs)
    ok("the title is in the document", "Report" in body)
    ok("the paragraph is in the document", "Hello world." in body)
    ok("bullets use the list style",
       any(p.style.name == "List Bullet" for p in d.paragraphs))
else:
    ok("docx refuses honestly", docbuilder.build("docx", {"sections": [{}]}, folder=OUT,
                                                 resolver=R).get("missing") == "python-docx")

# ════════════════════════════════════════════════════════════ SAFETY
sec("THE PATH JAIL HOLDS")
spec = {"title": "x", "sections": [{"heading": "h", "paragraphs": ["p"]}]}
for bad, why in [("/etc", "system folder"),
                 ("/", "filesystem root"),
                 (os.path.expanduser("~/.ssh"), "credential folder")]:
    r = docbuilder.build("docx", spec, folder=bad, resolver=R)
    ok(f"refuses to write into {why}", not r["ok"], r.get("message", "")[:60])

r = docbuilder.build("docx", {**spec, "filename": "../../escape"}, folder=OUT, resolver=R)
ok("a traversal filename cannot escape the folder",
   (not r["ok"]) or (os.path.realpath(OUT) in os.path.realpath(r.get("path", ""))),
   r.get("path") or r.get("message", "")[:60])

r = docbuilder.build("docx", {**spec, "filename": "evil.exe"}, folder=OUT, resolver=R)
ok("the extension is forced to match the type",
   r["ok"] and r["path"].endswith(".docx"), r.get("path", "")[-24:])

ok("an unknown type is refused, not guessed",
   not docbuilder.build("bat", spec, folder=OUT, resolver=R)["ok"])
ok("a non-dict spec is refused",
   not docbuilder.build("docx", "not a dict", folder=OUT, resolver=R)["ok"])
ok("an empty outline is refused",
   not docbuilder.build("docx", {"sections": []}, folder=OUT, resolver=R)["ok"])

if caps["pptx"]:
    big = {"title": "big", "slides": [{"title": f"s{i}"} for i in range(200)]}
    ok("an absurd slide count is capped, not attempted",
       not docbuilder.build("pptx", big, folder=OUT, resolver=R)["ok"])

sec("EXISTING FILES ARE NEVER SILENTLY OVERWRITTEN")
if caps["docx"]:
    a = docbuilder.build("docx", {**spec, "filename": "dup"}, folder=OUT, resolver=R)
    b = docbuilder.build("docx", {**spec, "filename": "dup"}, folder=OUT, resolver=R)
    ok("the second write gets its own name", a["path"] != b["path"],
       f"{os.path.basename(a['path'])} vs {os.path.basename(b['path'])}")
    ok("both files still exist",
       os.path.isfile(a["path"]) and os.path.isfile(b["path"]))

# ═══════════════════════════════════════════════════════════ ORGANISER
sec("ORGANISER: PREVIEW MOVES NOTHING")
D = os.path.join(os.path.expanduser("~"), "aura-test-org")
shutil.rmtree(D, ignore_errors=True)
os.makedirs(D)
FILES = ["a.jpg", "b.png", "r.pdf", "d.csv", "s.mp3", "i.exe",
         "c.py", "z.zip", "q.xyzzy", "n.txt"]
for n in FILES:
    open(os.path.join(D, n), "w").write("x" * 8)
os.makedirs(os.path.join(D, "Keep"), exist_ok=True)
open(os.path.join(D, ".hidden"), "w").write("x")
try:
    os.symlink("/etc/passwd", os.path.join(D, "evil.link"))
    HAVE_LINK = True
except Exception:
    HAVE_LINK = False

before = sorted(os.listdir(D))
p = organizer.plan(D, R)
ok("plan succeeds", p["ok"], p.get("message", "")[:60])
ok("NOTHING moved during the preview", sorted(os.listdir(D)) == before)
ok("every loose file is accounted for", p["total"] == len(FILES), str(p["total"]))
ok("a token is issued", bool(p.get("token")))
ok("images are grouped", p["counts"].get("Images") == 2, str(p["counts"]))
ok("unknown extensions go to Other", p["counts"].get("Other") == 1, str(p["counts"]))
skipped = {s["name"]: s["why"] for s in p["skipped"]}
ok("hidden files are skipped", skipped.get(".hidden") == "hidden", str(skipped))
ok("existing folders are skipped", skipped.get("Keep") == "folder", str(skipped))
if HAVE_LINK:
    ok("symlinks are never followed", skipped.get("evil.link") == "symlink", str(skipped))

sec("ORGANISER: NOTHING MOVES WITHOUT A VALID PLAN")
ok("a bogus token is refused", not organizer.apply(D, "not-a-token", R)["ok"])
ok("...and it says a preview is needed",
   organizer.apply(D, "nope", R).get("needsPlan") is True)
ok("still nothing moved", sorted(os.listdir(D)) == before)

stale = organizer.plan(D, R)
organizer._plans[stale["token"]]["at"] = time.time() - 9999
r = organizer.apply(D, stale["token"], R)
ok("an expired preview is refused", not r["ok"] and r.get("needsPlan"), r["message"][:50])

p2 = organizer.plan(D, R)
r = organizer.apply(os.path.expanduser("~"), p2["token"], R)
ok("a plan cannot be applied to a different folder", not r["ok"], r["message"][:50])

sec("ORGANISER: APPLY, RE-RUN, UNDO")
p3 = organizer.plan(D, R)
a = organizer.apply(D, p3["token"], R)
ok("apply succeeds", a["ok"], a.get("message", "")[:60])
ok("every file moved", a["moved"] == len(FILES), str(a["moved"]))
ok("no failures", not a["failed"], str(a["failed"]))
ok("Images/ holds the two images",
   sorted(os.listdir(os.path.join(D, "Images"))) == ["a.jpg", "b.png"])
ok("the untouched folder is still there", os.path.isdir(os.path.join(D, "Keep")))
ok("the hidden file was left alone", os.path.isfile(os.path.join(D, ".hidden")))
ok("a token cannot be replayed", not organizer.apply(D, p3["token"], R)["ok"])

again = organizer.plan(D, R)
ok("running it twice is a no-op", again["total"] == 0, str(again["total"]))

u = organizer.undo(D, R)
ok("undo succeeds", u["ok"], u.get("message", "")[:60])
ok("every file came back", u["restored"] == len(FILES), str(u["restored"]))
ok("the files are in the root again",
   all(os.path.isfile(os.path.join(D, n)) for n in FILES))
ok("empty category folders were cleaned up",
   not os.path.isdir(os.path.join(D, "Images")))

sec("ORGANISER: JAIL + CATEGORIES")
ok("refuses a folder outside the jail", not organizer.plan("/etc", R)["ok"])
ok("refuses a file instead of a folder",
   not organizer.plan(os.path.join(D, "a.jpg"), R)["ok"])
ok(".jpg -> Images", organizer.categorise("x.JPG") == "Images")
ok(".pptx -> Presentations", organizer.categorise("x.pptx") == "Presentations")
ok(".unknown -> Other", organizer.categorise("x.qqq") == "Other")
ok("no extension -> Other", organizer.categorise("README") == "Other")
oc = organizer.capabilities()
ok("categories are listed for the UI", len(oc["categories"]) >= 10, str(len(oc["categories"])))

# ═══════════════════════════════════════════════════════════════ QR
sec("QR PAIRING")
url = devices.pair_url(8000, "123456", host="192.168.1.50")
ok("the pair URL points at /phone", "/phone" in url, url)
ok("the code is embedded so scanning fills it in", "code=123456" in url, url)
q = devices.qr_svg(url)
if q.get("ok"):
    ok("a QR is produced", q["ok"])
    ok("it is an SVG", q["svg"].startswith("<svg") and q["svg"].endswith("</svg>"))
    ok("it declares a viewBox so it scales", "viewBox" in q["svg"])
    ok("it has a white background for scanners", 'fill="#ffffff"' in q["svg"])
    ok("it is compact (path, not thousands of rects)", len(q["svg"]) < 20000,
       f"{len(q['svg'])} bytes")
    # THE REAL TEST: does a decoder read it back?
    try:
        import re as _re, numpy as np, cv2
        n = q["modules"]
        img = np.ones((n, n), dtype=np.uint8) * 255
        for m in _re.finditer(r"M(\d+) (\d+)h(\d+)", q["svg"]):
            x, y, run = int(m.group(1)), int(m.group(2)), int(m.group(3))
            img[y, x:x + run] = 0
        big = cv2.resize(img, (n * 10, n * 10), interpolation=cv2.INTER_NEAREST)
        decoded, _, _ = cv2.QRCodeDetector().detectAndDecode(big)
        ok("A REAL DECODER READS IT BACK", decoded == url, decoded or "(no read)")
    except ImportError:
        print("  \033[33m~\033[0m opencv not installed — skipping the decode check")
    long_url = devices.pair_url(8000, "999999", host="192.168.100.200")
    ok("a longer URL still encodes", devices.qr_svg(long_url).get("ok"))
else:
    ok("QR absence is reported honestly, and the code still works",
       "qrcode" in q.get("message", "").lower(), q.get("message", "")[:60])

lan = devices.lan_ip()
ok("lan_ip() returns an address or None honestly",
   lan is None or (lan.count(".") == 3 and not lan.startswith("127.")), str(lan))
ok("pair_url falls back to localhost with no LAN",
   "localhost" in devices.pair_url(8000, "1", host=None) or lan is not None)

# ═══════════════════════════════════════════════════════ PPTX MEDIA
sec("PPTX MEDIA — IMAGES REALLY GET EMBEDDED")
if caps["pptx"]:
    import base64
    # 1×1 red PNG (valid, tiny).
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4z8DwHwAFBQIA"
        "X8jx0gAAAABJRU5ErkJggg==")
    img_path = os.path.join(OUT, "pixel.png")
    os.makedirs(OUT, exist_ok=True)
    with open(img_path, "wb") as fh:
        fh.write(png)
    r = docbuilder.build("pptx", {
        "title": "Media deck",
        "slides": [
            {"title": "Hero image", "kind": "image", "image": img_path,
             "imageCaption": "a generated pixel"},
            {"title": "Broken", "kind": "image",
             "image": "https://127.0.0.1:1/no-such.png"},
            {"title": "Text", "bullets": ["still renders"]},
        ]}, folder=OUT, resolver=R)
    ok("media deck still builds", r["ok"], r.get("message", "")[:70])
    names = zipfile.ZipFile(r["path"]).namelist()
    ok("the picture is inside ppt/media/",
       any(n.startswith("ppt/media/") for n in names), str(names[:8]))
    ok("one image embedded, one honest failure",
       r.get("embedded_images") == 1 and len(r.get("failed_images") or []) == 1,
       f"embedded={r.get('embedded_images')} failed={r.get('failed_images')}")
    ok("the failure names the slide and cause",
       bool(__import__("re").search(r"slide \d+", (r.get("failed_images") or [""])[0]))
       and "download" in (r.get("failed_images") or [""])[0],
       (r.get("failed_images") or [""])[0])
    # Pure resolver behaviour: local via jail, junk refused.
    p, o, n = docbuilder._load_image(img_path, R)
    ok("local path resolves through the jail resolver", o and os.path.isfile(p), n)
    p, o, n = docbuilder._load_image("/etc/passwd", R)
    ok("paths outside the home jail are refused", not o, n)
    p, o, n = docbuilder._load_image("", R)
    ok("empty source is refused, never embedded", not o, n)

shutil.rmtree(OUT, ignore_errors=True)
shutil.rmtree(D, ignore_errors=True)

print(f"\n  \033[32mPASS {len(P)}\033[0m  FAIL {len(F)}")
if F:
    print("  Failed: " + ", ".join(F))
    sys.exit(1)
